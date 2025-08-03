# Entry point for CLI or MCP tool integration
import argparse
import os
from datetime import datetime, timedelta
from parser import parse_transactions
from categorizer import categorize_transactions
from reporter import generate_report, generate_category_breakdown, filter_by_timeframe
from html_parser import extract_transactions_from_html
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import io
from flask import Flask, request, jsonify
from flask_cors import CORS
import json
from insights_generator import InsightsGenerator, FinancialProfile
from dotenv import load_dotenv
import seaborn as sns
import pandas as pd
import yagmail
from io import BytesIO
import textwrap
import numpy as np

load_dotenv()

# Initialize Flask app
app = Flask(__name__)
CORS(app)  # Enable CORS for all routes

def format_ai_text(text):
    """Format AI-generated text into bullet points and wrap long lines"""
    # Split into paragraphs
    paragraphs = text.split('\n')
    formatted = []
    
    for p in paragraphs:
        p = p.strip()
        if not p:
            continue
            
        # Convert to bullet point if not already
        if not p.startswith('•') and not p.startswith('-'):
            p = '• ' + p
            
        # Wrap long lines
        wrapped = textwrap.fill(p, width=80, subsequent_indent='  ')
        formatted.append(wrapped)
    
    return '\n'.join(formatted)

def apply_text_formatting(text_frame, font_size=11):
    """Apply consistent formatting to PowerPoint text frame"""
    from pptx.util import Pt
    
    # Clear any existing paragraphs
    text_frame.clear()
    
    # Add a new paragraph with the specified font size
    p = text_frame.add_paragraph()
    p.font.size = Pt(font_size)
    
    return p

def create_daily_spending_chart(transactions, start_date, end_date):
    """Create a chart showing daily spending trends"""
    # Convert dates to datetime
    start = datetime.strptime(start_date, "%Y-%m-%d")
    end = datetime.strptime(end_date, "%Y-%m-%d")
    
    # Create date range
    date_range = pd.date_range(start=start, end=end, freq='D')
    daily_spend = {date: 0 for date in date_range}
    
    # Sum transactions by date
    for tx in transactions:
        tx_date = tx['timestamp'] if isinstance(tx['timestamp'], datetime) else datetime.strptime(tx['timestamp'].split()[0], "%Y-%m-%d")
        if start <= tx_date <= end:
            daily_spend[tx_date.replace(hour=0, minute=0, second=0, microsecond=0)] += float(tx['amount'])
    
    # Create DataFrame
    df = pd.DataFrame(list(daily_spend.items()), columns=['date', 'amount'])
    
    # Create figure
    plt.figure(figsize=(12, 6))
    sns.lineplot(data=df, x='date', y='amount', marker='o')
    plt.title('Daily Spending Trend')
    plt.xlabel('Date')
    plt.ylabel('Amount (₹)')
    plt.xticks(rotation=45)
    plt.grid(True)
    
    # Save to BytesIO
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_category_chart(transactions):
    """Create a chart showing spending by category"""
    # Sum transactions by category
    category_spend = {}
    for tx in transactions:
        category = tx['category']
        if category not in category_spend:
            category_spend[category] = 0
        category_spend[category] += float(tx['amount'])
    
    # Create DataFrame
    df = pd.DataFrame(list(category_spend.items()), columns=['category', 'amount'])
    df = df.sort_values('amount', ascending=True)
    
    # Create figure
    plt.figure(figsize=(10, 8))
    sns.barplot(data=df, x='amount', y='category')
    plt.title('Spending by Category')
    plt.xlabel('Amount (₹)')
    plt.ylabel('Category')
    plt.grid(True)
    
    # Save to BytesIO
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_category_trend_chart(transactions, start_date, end_date):
    """Create a chart showing category trends over time"""
    # Convert dates to datetime
    start = datetime.strptime(start_date, "%Y-%m-%d")
    end = datetime.strptime(end_date, "%Y-%m-%d")
    
    # Create weekly bins
    date_range = pd.date_range(start=start, end=end, freq='W')
    category_trends = {date: {} for date in date_range}
    
    # Sum transactions by week and category
    for tx in transactions:
        tx_date = tx['timestamp'] if isinstance(tx['timestamp'], datetime) else datetime.strptime(tx['timestamp'].split()[0], "%Y-%m-%d")
        if start <= tx_date <= end:
            # Find the corresponding week
            week = date_range[date_range.searchsorted(tx_date) - 1]
            category = tx['category']
            if category not in category_trends[week]:
                category_trends[week][category] = 0
            category_trends[week][category] += float(tx['amount'])
    
    # Create DataFrame
    data = []
    for week, categories in category_trends.items():
        for category, amount in categories.items():
            data.append({'week': week, 'category': category, 'amount': amount})
    df = pd.DataFrame(data)
    
    # Create figure
    plt.figure(figsize=(15, 8))
    sns.lineplot(data=df, x='week', y='amount', hue='category')
    plt.title('Category Spending Trends')
    plt.xlabel('Week')
    plt.ylabel('Amount (₹)')
    plt.xticks(rotation=45)
    plt.legend(bbox_to_anchor=(1.05, 1), loc='upper left')
    plt.grid(True)
    
    # Save to BytesIO
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_weekday_spending_chart(transactions):
    """Create a chart showing spending patterns by day of week"""
    weekday_spend = {i: 0 for i in range(7)}  # 0 = Monday, 6 = Sunday
    weekday_count = {i: 0 for i in range(7)}
    
    for tx in transactions:
        tx_date = tx['timestamp'] if isinstance(tx['timestamp'], datetime) else datetime.strptime(tx['timestamp'].split()[0], "%Y-%m-%d")
        weekday = tx_date.weekday()
        weekday_spend[weekday] += float(tx['amount'])
        weekday_count[weekday] += 1
    
    # Calculate average spend per day
    weekday_avg = {day: (spend/weekday_count[day] if weekday_count[day] > 0 else 0) 
                  for day, spend in weekday_spend.items()}
    
    # Create DataFrame
    days = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
    df = pd.DataFrame({
        'day': days,
        'total_spend': [weekday_spend[i] for i in range(7)],
        'avg_spend': [weekday_avg[i] for i in range(7)],
        'transaction_count': [weekday_count[i] for i in range(7)]
    })
    
    # Create figure with two subplots side by side
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(15, 6))
    
    # Total spend by day
    sns.barplot(data=df, x='day', y='total_spend', ax=ax1, color='skyblue')
    ax1.set_title('Total Spending by Day of Week')
    ax1.set_ylabel('Total Amount (₹)')
    ax1.tick_params(axis='x', rotation=45)
    
    # Average spend by day
    sns.barplot(data=df, x='day', y='avg_spend', ax=ax2, color='lightgreen')
    ax2.set_title('Average Transaction Size by Day of Week')
    ax2.set_ylabel('Average Amount (₹)')
    ax2.tick_params(axis='x', rotation=45)
    
    plt.tight_layout(pad=2.0)  # Increase padding
    
    # Save to BytesIO with higher quality
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_time_of_day_chart(transactions):
    """Create a chart showing spending patterns by time of day"""
    # Define time periods
    time_periods = {
        'Early Morning (12AM-6AM)': (0, 6),
        'Morning (6AM-12PM)': (6, 12),
        'Afternoon (12PM-6PM)': (12, 18),
        'Evening (6PM-12AM)': (18, 24)
    }
    
    # Initialize counters
    period_spend = {period: 0 for period in time_periods}
    period_count = {period: 0 for period in time_periods}
    
    for tx in transactions:
        if isinstance(tx['timestamp'], str):
            try:
                time = datetime.strptime(tx['timestamp'], "%Y-%m-%d %H:%M:%S").hour
            except:
                time = datetime.strptime(tx['timestamp'], "%Y-%m-%d").hour
        else:
            time = tx['timestamp'].hour
            
        # Find which period this hour belongs to
        for period, (start, end) in time_periods.items():
            if start <= time < end:
                period_spend[period] += float(tx['amount'])
                period_count[period] += 1
                break
    
    # Calculate averages
    period_avg = {period: (spend/period_count[period] if period_count[period] > 0 else 0) 
                 for period, spend in period_spend.items()}
    
    # Create DataFrame
    df = pd.DataFrame({
        'period': list(time_periods.keys()),
        'total_spend': list(period_spend.values()),
        'avg_spend': list(period_avg.values()),
        'transaction_count': list(period_count.values())
    })
    
    # Create figure with two subplots
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(15, 6))
    
    # Total spend by time period
    sns.barplot(data=df, x='period', y='total_spend', ax=ax1, color='purple')
    ax1.set_title('Total Spending by Time of Day')
    ax1.set_ylabel('Total Amount (₹)')
    ax1.tick_params(axis='x', rotation=45)
    
    # Transaction count by time period
    sns.barplot(data=df, x='period', y='transaction_count', ax=ax2, color='orange')
    ax2.set_title('Number of Transactions by Time of Day')
    ax2.set_ylabel('Number of Transactions')
    ax2.tick_params(axis='x', rotation=45)
    
    plt.tight_layout()
    
    # Save to BytesIO
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_transaction_size_chart(transactions):
    """Create a chart showing transaction size distribution"""
    amounts = [float(tx['amount']) for tx in transactions]
    
    # Create figure with two subplots
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(15, 6))
    
    # Histogram of transaction sizes
    sns.histplot(amounts, bins=30, ax=ax1, color='teal')
    ax1.set_title('Distribution of Transaction Sizes')
    ax1.set_xlabel('Amount (₹)')
    ax1.set_ylabel('Number of Transactions')
    
    # Box plot of transaction sizes by category
    df = pd.DataFrame([(tx['category'], float(tx['amount'])) for tx in transactions],
                     columns=['category', 'amount'])
    sns.boxplot(data=df, x='category', y='amount', ax=ax2)
    ax2.set_title('Transaction Sizes by Category')
    ax2.set_xlabel('Category')
    ax2.set_ylabel('Amount (₹)')
    ax2.tick_params(axis='x', rotation=45)
    
    plt.tight_layout()
    
    # Save to BytesIO
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_top_categories_chart(transactions):
    """Create a detailed chart of top 3 spending categories"""
    # Calculate spending by category
    category_spend = {}
    for tx in transactions:
        category = tx['category']
        if category not in category_spend:
            category_spend[category] = {
                'total': 0,
                'transactions': [],
                'count': 0
            }
        amount = float(tx['amount'])
        category_spend[category]['total'] += amount
        category_spend[category]['transactions'].append(amount)
        category_spend[category]['count'] += 1

    # Get top 3 categories
    top_categories = sorted(category_spend.items(), key=lambda x: x[1]['total'], reverse=True)[:3]
    
    # Create figure with multiple subplots
    fig = plt.figure(figsize=(15, 8))
    gs = fig.add_gridspec(2, 3)
    
    # Total amounts (bar chart)
    ax1 = fig.add_subplot(gs[0, :])
    categories = [cat[0] for cat in top_categories]
    amounts = [cat[1]['total'] for cat in top_categories]
    bars = ax1.bar(categories, amounts)
    ax1.set_title('Top 3 Categories by Total Spend')
    ax1.set_ylabel('Total Amount (₹)')
    
    # Add value labels on bars
    for bar in bars:
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height,
                f'₹{height:,.0f}',
                ha='center', va='bottom')
    
    # Transaction count comparison
    ax2 = fig.add_subplot(gs[1, 0])
    counts = [cat[1]['count'] for cat in top_categories]
    ax2.pie(counts, labels=categories, autopct='%1.1f%%')
    ax2.set_title('Transaction Count Distribution')
    
    # Average transaction size
    ax3 = fig.add_subplot(gs[1, 1])
    avgs = [cat[1]['total']/cat[1]['count'] for cat in top_categories]
    ax3.bar(categories, avgs)
    ax3.set_title('Average Transaction Size')
    ax3.set_ylabel('Amount (₹)')
    
    # Monthly trend
    ax4 = fig.add_subplot(gs[1, 2])
    for cat, data in top_categories:
        amounts = sorted(data['transactions'])
        ax4.hist(amounts, bins=20, alpha=0.3, label=cat)
    ax4.set_title('Transaction Size Distribution')
    ax4.set_xlabel('Amount (₹)')
    ax4.legend()
    
    plt.tight_layout(pad=2.0)
    
    # Save to BytesIO
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_financial_distribution_chart(transactions):
    """Create a chart showing spend vs investment vs savings"""
    # Define category types
    investment_categories = {'Investment', 'Stocks', 'Mutual Funds', 'FD', 'Fixed Deposit'}
    expense_categories = {'Food', 'Transport', 'Shopping', 'Bills', 'Groceries', 'Entertainment', 'Outing'}
    
    # Calculate totals
    total_investment = sum(float(tx['amount']) for tx in transactions 
                         if tx['category'] in investment_categories)
    total_expenses = sum(float(tx['amount']) for tx in transactions 
                        if tx['category'] in expense_categories)
    total_amount = sum(float(tx['amount']) for tx in transactions)
    total_other = total_amount - (total_investment + total_expenses)  # Potential savings
    
    # Create figure with subplots
    fig = plt.figure(figsize=(15, 6))
    
    # Pie chart
    ax1 = plt.subplot(121)
    labels = ['Expenses', 'Investments', 'Other/Savings']
    sizes = [total_expenses, total_investment, total_other]
    colors = ['#ff9999', '#66b3ff', '#99ff99']
    
    patches, texts, autotexts = ax1.pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%',
                                       startangle=90)
    ax1.axis('equal')
    ax1.set_title('Financial Distribution')
    
    # Add legend with amounts
    legend_labels = [
        f'Expenses: ₹{total_expenses:,.2f}',
        f'Investments: ₹{total_investment:,.2f}',
        f'Other/Savings: ₹{total_other:,.2f}'
    ]
    ax1.legend(patches, legend_labels, title="Breakdown", loc="center left", bbox_to_anchor=(1, 0, 0.5, 1))
    
    # Bar chart showing monthly progression
    ax2 = plt.subplot(122)
    monthly_data = {}
    
    for tx in transactions:
        date = tx['timestamp'] if isinstance(tx['timestamp'], datetime) else datetime.strptime(tx['timestamp'].split()[0], "%Y-%m-%d")
        month_key = date.strftime("%Y-%m")
        if month_key not in monthly_data:
            monthly_data[month_key] = {'investment': 0, 'expenses': 0, 'other': 0}
        
        amount = float(tx['amount'])
        if tx['category'] in investment_categories:
            monthly_data[month_key]['investment'] += amount
        elif tx['category'] in expense_categories:
            monthly_data[month_key]['expenses'] += amount
        else:
            monthly_data[month_key]['other'] += amount
    
    months = sorted(monthly_data.keys())
    investments = [monthly_data[m]['investment'] for m in months]
    expenses = [monthly_data[m]['expenses'] for m in months]
    others = [monthly_data[m]['other'] for m in months]
    
    x = range(len(months))
    width = 0.25
    
    ax2.bar([i - width for i in x], expenses, width, label='Expenses', color='#ff9999')
    ax2.bar(x, investments, width, label='Investments', color='#66b3ff')
    ax2.bar([i + width for i in x], others, width, label='Other/Savings', color='#99ff99')
    
    ax2.set_title('Monthly Distribution')
    ax2.set_xlabel('Month')
    ax2.set_ylabel('Amount (₹)')
    ax2.set_xticks(x)
    ax2.set_xticklabels(months, rotation=45)
    ax2.legend()
    
    plt.tight_layout(pad=2.0)
    
    # Save to BytesIO
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_category_detailed_chart(transactions, category_name):
    """Create a detailed analysis of a specific spending category"""
    # Filter transactions for the category
    category_transactions = [tx for tx in transactions 
                           if tx['category'] == category_name]
    
    if not category_transactions:
        return None
    
    # Create figure with subplots
    fig = plt.figure(figsize=(15, 10))
    gs = fig.add_gridspec(2, 2)
    
    # Daily spending pattern
    ax1 = fig.add_subplot(gs[0, :])
    dates = [tx['timestamp'] if isinstance(tx['timestamp'], datetime) 
            else datetime.strptime(tx['timestamp'].split()[0], "%Y-%m-%d")
            for tx in category_transactions]
    amounts = [float(tx['amount']) for tx in category_transactions]
    
    ax1.plot(dates, amounts, marker='o', linestyle='-', alpha=0.6)
    ax1.set_title(f'{category_name} Spending Pattern')
    ax1.set_xlabel('Date')
    ax1.set_ylabel('Amount (₹)')
    ax1.tick_params(axis='x', rotation=45)
    
    # Add total spent annotation
    total_spent = sum(amounts)
    avg_per_tx = total_spent / len(amounts)
    ax1.text(0.02, 0.98, 
             f'Total Spent: ₹{total_spent:,.2f}\nAvg per Transaction: ₹{avg_per_tx:,.2f}',
             transform=ax1.transAxes,
             verticalalignment='top',
             bbox=dict(boxstyle='round', facecolor='white', alpha=0.8))
    
    # Transaction size distribution
    ax2 = fig.add_subplot(gs[1, 0])
    n, bins, patches = ax2.hist(amounts, bins=20, color='purple', alpha=0.6)
    ax2.set_title(f'{category_name} Transaction Size Distribution')
    ax2.set_xlabel('Amount (₹)')
    ax2.set_ylabel('Number of Transactions')
    
    # Add median and mean lines
    median = np.median(amounts)
    mean = np.mean(amounts)
    ax2.axvline(median, color='red', linestyle='dashed', alpha=0.8, label=f'Median: ₹{median:,.0f}')
    ax2.axvline(mean, color='green', linestyle='dashed', alpha=0.8, label=f'Mean: ₹{mean:,.0f}')
    ax2.legend()
    
    # Day of week analysis
    ax3 = fig.add_subplot(gs[1, 1])
    day_spending = {i: [] for i in range(7)}  # 0 = Monday, 6 = Sunday
    for date, amount in zip(dates, amounts):
        day_spending[date.weekday()].append(amount)
    
    days = ['Mon', 'Tue', 'Wed', 'Thu', 'Fri', 'Sat', 'Sun']
    avg_by_day = [sum(amounts)/len(amounts) if amounts else 0 
                  for amounts in day_spending.values()]
    
    bars = ax3.bar(days, avg_by_day)
    ax3.set_title(f'{category_name} Average Spending by Day')
    ax3.set_ylabel('Average Amount (₹)')
    
    # Add value labels on bars
    for bar in bars:
        height = bar.get_height()
        ax3.text(bar.get_x() + bar.get_width()/2., height,
                f'₹{height:,.0f}',
                ha='center', va='bottom')
    
    plt.tight_layout(pad=2.0)
    
    # Save to BytesIO
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_ppt(start_date: str, end_date: str, output: str, email: str = None, include_insights: bool = False):
    """Create PowerPoint presentation with financial insights"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import textwrap
    
    # Get transactions
    transactions = parse_transactions("data/transactions.csv")
    
    # Filter transactions by date range
    start = datetime.strptime(start_date, "%Y-%m-%d")
    end = datetime.strptime(end_date, "%Y-%m-%d")
    filtered_transactions = [
        tx for tx in transactions 
        if start <= (tx['timestamp'] if isinstance(tx['timestamp'], datetime) else datetime.strptime(tx['timestamp'].split()[0], "%Y-%m-%d")) <= end
    ]

    # Categorize transactions
    print("🤖 Categorizing transactions...")
    try:
        categorized_transactions = categorize_transactions(filtered_transactions)
        print(f"✅ Categorized {len(categorized_transactions)} transactions")
    except Exception as e:
        print(f"❌ Error categorizing transactions: {e}")
        return

    # Generate AI insights if requested
    insights = None
    if include_insights:
        print("\n🤖 Generating AI insights...")
        try:
            # Load profile data
            with open("profile.json") as f:
                profile_data = json.load(f)
            
            profile = FinancialProfile(
                monthly_income=profile_data['monthly_income'],
                monthly_rent=profile_data['monthly_rent'],
                weekly_travel=profile_data['weekly_travel'],
                weekly_groceries=profile_data['weekly_groceries'],
                city=profile_data['city'],
                lifestyle=profile_data['lifestyle'],
                relationship_status=profile_data['relationship_status']
            )
            
            generator = InsightsGenerator()
            insights = generator.generate_insights(profile, categorized_transactions, start_date, end_date)
            print("✅ Generated AI insights")

        except Exception as e:
            print(f"❌ Error generating insights: {e}")
            insights = None
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Financial Insights Report"
    subtitle.text = f"Period: {start_date} to {end_date}"

    # Daily spending trend slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Daily Spending Trend"
    daily_chart = create_daily_spending_chart(categorized_transactions, start_date, end_date)
    pic = slide.shapes.add_picture(daily_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Category breakdown slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Spending by Category"
    category_chart = create_category_chart(categorized_transactions)
    pic = slide.shapes.add_picture(category_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Category trends slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Category Spending Trends"
    trend_chart = create_category_trend_chart(categorized_transactions, start_date, end_date)
    pic = slide.shapes.add_picture(trend_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Day of week analysis slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Spending Patterns by Day of Week"
    weekday_chart = create_weekday_spending_chart(categorized_transactions)
    pic = slide.shapes.add_picture(weekday_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Time of day analysis slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Spending Patterns by Time of Day"
    time_chart = create_time_of_day_chart(categorized_transactions)
    pic = slide.shapes.add_picture(time_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Transaction size analysis slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Transaction Size Analysis"
    size_chart = create_transaction_size_chart(categorized_transactions)
    pic = slide.shapes.add_picture(size_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Top 3 Categories Analysis slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Top 3 Categories Deep Dive"
    top_cats_chart = create_top_categories_chart(categorized_transactions)
    pic = slide.shapes.add_picture(top_cats_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Financial Distribution slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Spend vs Investment vs Savings"
    dist_chart = create_financial_distribution_chart(categorized_transactions)
    pic = slide.shapes.add_picture(dist_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Entertainment Analysis slide
    ent_chart = create_category_detailed_chart(categorized_transactions, "Entertainment")
    if ent_chart:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Entertainment Spending Analysis"
        pic = slide.shapes.add_picture(ent_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Outing Analysis slide
    outing_chart = create_category_detailed_chart(categorized_transactions, "Outing")
    if outing_chart:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Outing Spending Analysis"
        pic = slide.shapes.add_picture(outing_chart, Inches(0.5), Inches(1.5), width=Inches(9))

    # Summary statistics slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Summary Statistics"
    
    total_spent = sum(float(tx['amount']) for tx in categorized_transactions)
    avg_daily = total_spent / (end - start).days if (end - start).days > 0 else total_spent
    
    tf = body.text_frame
    tf.text = f"Total Spent: ₹{total_spent:,.2f}\n"
    tf.add_paragraph().text = f"Average Daily Spend: ₹{avg_daily:,.2f}"
    tf.add_paragraph().text = f"Number of Transactions: {len(categorized_transactions)}"

    # AI Insights slides (after all charts)
    if insights and 'insights' in insights and include_insights:
        # Overview and Pattern Analysis
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "AI Analysis: Spending Patterns"
        
        # Format text and add content
        body.text_frame.clear()  # Clear existing content
        for line in format_ai_text(insights['insights']['overview']).split('\n'):
            p = body.text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
        
        # Areas for Improvement
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "AI Analysis: Areas for Improvement"
        
        # Format text and add content
        body.text_frame.clear()  # Clear existing content
        for line in format_ai_text(insights['insights']['improvements']).split('\n'):
            p = body.text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
        
        # Lifestyle Analysis
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "AI Analysis: Lifestyle Impact"
        
        # Format text and add content
        body.text_frame.clear()  # Clear existing content
        for line in format_ai_text(insights['insights']['lifestyle']).split('\n'):
            p = body.text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
        
        # Recommendations
        recommendations = format_ai_text(insights['insights']['recommendations'])
        rec_lines = recommendations.split('\n')
        
        if len(rec_lines) > 10:  # Split into two slides if too long
            # First recommendations slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = "AI Recommendations"
            
            # Format text and add first half of content
            body.text_frame.clear()  # Clear existing content
            for line in rec_lines[:10]:
                p = body.text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(11)
            
            # Second recommendations slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = "AI Recommendations (continued)"
            
            # Format text and add second half of content
            body.text_frame.clear()  # Clear existing content
            for line in rec_lines[10:]:
                p = body.text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(11)
        else:
            # Single recommendations slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = "AI Recommendations"
            
            # Format text and add content
            body.text_frame.clear()  # Clear existing content
            for line in rec_lines:
                p = body.text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(11)
    
    # Save presentation
    prs.save(output)
    print(f"\n✅ PowerPoint saved to {output}")
    
    # Send email if requested
    if email:
        try:
            # Initialize yagmail with your Gmail credentials
            gmail_user = os.getenv('GMAIL_USER')
            gmail_pass = os.getenv('GMAIL_APP_PASSWORD')
            
            if not gmail_user or not gmail_pass:
                print("❌ Gmail credentials not found in environment variables")
                return
            
            yag = yagmail.SMTP(gmail_user, gmail_pass)
            
            # Prepare email content
            subject = f"Financial Insights Report ({start_date} to {end_date})"
            contents = [
                "Please find attached your financial insights report.",
                f"Period: {start_date} to {end_date}",
                f"Total Spent: ₹{total_spent:,.2f}",
                f"Average Daily Spend: ₹{avg_daily:,.2f}",
                output  # Attach the PowerPoint file
            ]
            
            # Send email
            yag.send(email, subject, contents)
            print(f"✅ Report sent to {email}")
            
        except Exception as e:
            print(f"❌ Error sending email: {e}")

def daily_range(start_date, end_date):
    """Generate daily reports for each day between start_date and end_date"""
    print("Reading transactions from CSV (Google Takeout export)...")
    transactions = parse_transactions("data/transactions.csv")
    print(f"Parsed {len(transactions)} transactions.")

    # Convert dates to datetime objects
    start = datetime.strptime(start_date, "%Y-%m-%d")
    end = datetime.strptime(end_date, "%Y-%m-%d")
    
    # Initialize results array
    daily_reports = []
    
    # For each day in the range
    current = start
    while current <= end:
        # Filter transactions for just this day
        day_start = current.replace(hour=0, minute=0, second=0)
        day_end = (current + timedelta(days=1)).replace(hour=0, minute=0, second=0)
        day_transactions = [
            tx for tx in transactions 
            if day_start <= (tx['timestamp'] if isinstance(tx['timestamp'], datetime) else datetime.strptime(tx['timestamp'].split()[0], "%Y-%m-%d")) < day_end
        ]
        
        if day_transactions:
            # Categorize transactions for this day
            categorized = categorize_transactions(day_transactions)
            
            # Calculate totals by category
            categories = {}
            for tx in categorized:
                category = tx['category']
                amount = tx['amount']
                if category not in categories:
                    categories[category] = {
                        'total': 0,
                        'count': 0,
                        'transactions': []
                    }
                categories[category]['total'] += amount
                categories[category]['count'] += 1
                categories[category]['transactions'].append({
                    'description': tx['description'],
                    'amount': tx['amount'],
                    'merchant': tx['merchant'],
                    'timestamp': tx['timestamp'].strftime("%Y-%m-%d %H:%M:%S")  # Convert to string for JSON
                })
            
            # Create daily summary
            daily_summary = {
                'date': current.strftime("%Y-%m-%d"),
                'total_spent': sum(cat['total'] for cat in categories.values()),
                'total_transactions': sum(cat['count'] for cat in categories.values()),
                'categories': [
                    {
                        'name': category,
                        'total': data['total'],
                        'count': data['count'],
                        'transactions': data['transactions']
                    }
                    for category, data in sorted(
                        categories.items(),
                        key=lambda x: x[1]['total'],
                        reverse=True
                    )
                ]
            }
            
            daily_reports.append(daily_summary)
        else:
            # Add empty summary for days with no transactions
            daily_reports.append({
                'date': current.strftime("%Y-%m-%d"),
                'total_spent': 0,
                'total_transactions': 0,
                'categories': []
            })
        
        current += timedelta(days=1)
    
    # Print summary
    print(f"\nGenerated reports for {len(daily_reports)} days from {start_date} to {end_date}")
    print(f"Found transactions on {len([d for d in daily_reports if d['total_transactions'] > 0])} days")
    
    # Print each day's summary
    for day in daily_reports:
        if day['total_transactions'] > 0:
            print(f"\n📅 {day['date']}")
            print(f"💰 Total: ₹{day['total_spent']:.2f} ({day['total_transactions']} transactions)")
            for cat in day['categories']:
                print(f"  {cat['name']}: ₹{cat['total']:.2f} ({cat['count']} txns)")
    
    return daily_reports

def run(mode):
    print("Reading transactions from CSV (Google Takeout export)...")
    transactions = parse_transactions("data/transactions.csv")
    print(f"Parsed {len(transactions)} transactions.")

    # Filter transactions by timeframe first to speed up categorization
    transactions = filter_by_timeframe(transactions, mode)
    print(f"Found {len(transactions)} transactions for {mode} period.")
    
    if not transactions:
        print(f"No transactions found for {mode} period.")
        return

    categorized = categorize_transactions(transactions)

    print(f"Generating {mode} report...")
    report = generate_report(categorized, mode)
    print(report)

def breakdown(mode, category=None):
    print("Reading transactions from CSV (Google Takeout export)...")
    transactions = parse_transactions("data/transactions.csv")
    print(f"Parsed {len(transactions)} transactions.")

    # Filter transactions by timeframe first to speed up categorization
    transactions = filter_by_timeframe(transactions, mode)
    print(f"Found {len(transactions)} transactions for {mode} period.")
    
    if not transactions:
        print(f"No transactions found for {mode} period.")
        return

    print("🤖 Categorizing transactions with improved ML model...")
    categorized = categorize_transactions(transactions)

    print(f"Generating {mode} breakdown...")
    breakdown_report = generate_category_breakdown(categorized, mode, category)
    print(breakdown_report)

def setup():
    """Initial setup: Convert HTML to CSV and train ML model"""
    print("🔄 Step 1: Converting HTML to CSV...")
    try:
        from html_parser import extract_transactions_from_html
        html_path = "data/MyActivity.html"
        csv_path = "data/transactions.csv"
        extract_transactions_from_html(html_path, csv_path)
        print("✅ Transactions extracted and saved to data/transactions.csv")
        
    except Exception as e:
        print(f"❌ Error extracting transactions: {e}")
        return
    
    print("\n🤖 Step 2: Training improved ML model...")
    try:
        # Import and run the training script
        import sys
        import os
        
        # Add model directory to path
        model_dir = os.path.join(os.path.dirname(__file__), 'model')
        if model_dir not in sys.path:
            sys.path.append(model_dir)
            
        from train_model_improved import train_improved_model
        success = train_improved_model()
        
        if success:
            print("✅ Setup completed successfully!")
        else:
            print("✅ Setup completed (HTML extracted, ML training failed)")
            print("   To enable ML model, add labeled data to model/transaction_categorization_training_data.csv")
            
    except Exception as e:
        print(f"❌ Error training ML model: {e}")
        print("   The system will use LLM + rule-based categorization as fallback.")
        print("✅ Setup completed (HTML extracted, ML training failed)")

def generate_insights(start_date: str, end_date: str, profile_file: str):
    """Generate AI-powered insights for the given period"""
    # Load profile data
    with open(profile_file) as f:
        profile_data = json.load(f)
    
    # Create profile
    profile = FinancialProfile(
        monthly_income=profile_data['monthly_income'],
        monthly_rent=profile_data['monthly_rent'],
        weekly_travel=profile_data['weekly_travel'],
        weekly_groceries=profile_data['weekly_groceries'],
        city=profile_data['city'],
        lifestyle=profile_data['lifestyle'],
        relationship_status=profile_data['relationship_status']
    )
    
    # Get transactions
    transactions = parse_transactions("data/transactions.csv")
    
    # Generate insights
    generator = InsightsGenerator()
    insights = generator.generate_insights(profile, transactions, start_date, end_date)
    
    if "error" in insights:
        print(f"\n❌ Error: {insights['error']}")
        if "details" in insights:
            print(f"Details: {insights['details']}")
        return
    
    # Print insights
    print("\n=== Financial Insights Report ===")
    print(f"Period: {start_date} to {end_date}")
    print("\n=== Overview ===")
    print(insights['insights']['overview'])
    print("\n=== Areas for Improvement ===")
    print(insights['insights']['improvements'])
    print("\n=== Lifestyle Analysis ===")
    print(insights['insights']['lifestyle'])
    print("\n=== Recommendations ===")
    print(insights['insights']['recommendations'])
    
    # Save insights to file
    output_file = f"insights_{start_date}_{end_date}.json"
    with open(output_file, 'w') as f:
        json.dump(insights, f, indent=2)
    print(f"\n✅ Detailed insights saved to {output_file}")

# Add Flask routes
@app.route('/api/chart-data')
def get_chart_data():
    mode = request.args.get('mode', 'monthly')
    if mode not in ['daily', 'weekly', 'monthly', 'quarterly']:
        return jsonify({"error": "Invalid mode"}), 400

    try:
        transactions = parse_transactions("data/transactions.csv")
        transactions = filter_by_timeframe(transactions, mode)
        
        if not transactions:
            return jsonify({"error": f"No transactions found for {mode} period."})

        categorized = categorize_transactions(transactions)
        
        # Convert the report data into a format suitable for charts
        categories = {}
        for tx in categorized:
            category = tx['category']
            amount = tx['amount']
            if category not in categories:
                categories[category] = {
                    'total': 0,
                    'count': 0
                }
            categories[category]['total'] += amount
            categories[category]['count'] += 1

        # Sort by total amount
        sorted_data = sorted(
            [{"category": k, "amount": v['total'], "count": v['count']} 
             for k, v in categories.items()],
            key=lambda x: x['amount'],
            reverse=True
        )

        return jsonify({
            "data": sorted_data,
            "mode": mode,
            "total_transactions": len(transactions)
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/daily-range')
def get_daily_range_data():
    start_date = request.args.get('start_date')
    end_date = request.args.get('end_date')
    
    if not start_date or not end_date:
        return jsonify({"error": "start_date and end_date are required"}), 400
        
    try:
        reports = daily_range(start_date, end_date)
        return jsonify(reports)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    # Check if running as Flask app
    if os.environ.get('FLASK_RUN'):
        app.run(port=8000)
    else:
        # Original CLI logic
        parser = argparse.ArgumentParser()
        parser.add_argument("command", choices=["setup", "run", "breakdown", "daily_range", "create_ppt", "insights"], 
                        help="setup = convert HTML to CSV, run = generate report, breakdown = detailed category analysis, daily_range = daily reports between dates, create_ppt = create PowerPoint presentation, insights = generate AI insights")
        parser.add_argument("--mode", choices=["daily", "weekly", "monthly", "quarterly"], default="weekly")
        parser.add_argument("--category", type=str, help="Specific category to analyze (for breakdown command)")
        parser.add_argument("--start-date", type=str, help="Start date in YYYY-MM-DD format (for daily_range and create_ppt commands)")
        parser.add_argument("--end-date", type=str, help="End date in YYYY-MM-DD format (for daily_range and create_ppt commands)")
        parser.add_argument("--output", type=str, default="financial_insights.pptx", help="Output PowerPoint file name (for create_ppt command)")
        parser.add_argument("--email", type=str, help="Email address to send the report to (for create_ppt command)")
        parser.add_argument("--profile", type=str, default="profile.json", help="Profile JSON file (for insights command)")
        parser.add_argument("--insights", action="store_true", help="Include AI-generated insights in the PowerPoint")
        args = parser.parse_args()

        if args.command == "setup":
            setup()
        elif args.command == "run":
            run(args.mode)
        elif args.command == "breakdown":
            breakdown(args.mode, args.category)
        elif args.command == "daily_range":
            if not args.start_date or not args.end_date:
                print("Error: --start-date and --end-date are required for daily_range command")
                parser.print_help()
                exit(1)
            daily_range(args.start_date, args.end_date)
        elif args.command == "create_ppt":
            if not args.start_date or not args.end_date:
                print("Error: --start-date and --end-date are required for create_ppt command")
                parser.print_help()
                exit(1)
            create_ppt(args.start_date, args.end_date, args.output, args.email, args.insights)
        elif args.command == "insights":
            if not args.start_date or not args.end_date:
                print("Error: --start-date and --end-date are required for insights command")
                parser.print_help()
                exit(1)
            generate_insights(args.start_date, args.end_date, args.profile)

# Usage examples:
# python main.py setup                     # Extract transactions + train ML model
# python main.py run --mode daily          # Generate daily report with ML model
# python main.py run --mode weekly         # Generate weekly report with ML model  
# python main.py run --mode monthly        # Generate monthly report with ML model
# python main.py run --mode quarterly      # Generate quarterly report with ML model
# python main.py breakdown --mode monthly  # Show category breakdown with transaction counts
# python main.py breakdown --mode monthly --category Outing  # Detailed Outing transactions
# python main.py daily_range --start-date 2024-01-01 --end-date 2024-01-31  # Daily reports for January 2024
# python main.py create_ppt --start-date 2025-01-01 --end-date 2025-01-31 --output report.pptx --email banerjeeaviroop01@gmail.com --insights

# To run as API server:
# FLASK_RUN=1 python main.py
